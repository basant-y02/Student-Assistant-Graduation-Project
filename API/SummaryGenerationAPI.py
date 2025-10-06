import os
import subprocess
from flask import Flask, request, jsonify
import assemblyai as aai
import yt_dlp as youtube_dl
import json
from flask_cors import CORS
import google.generativeai as genai
from werkzeug.utils import secure_filename
import mysql.connector
from PyPDF2 import PdfReader
from pptx import Presentation
from docx import Document
from google.generativeai.types import HarmCategory, HarmBlockThreshold
import re

global AIResponce
AIResponce = ""

def convert_text_to_html(text):
    # Convert headers
    text = re.sub(r'## (.*?)\n', r'<h2>\1</h2>\n', text)
    text = re.sub(r'# (.*?)\n', r'<h1>\1</h1>\n', text)

    # Convert bold text
    text = re.sub(r'\*\*(.*?)\*\*', r'<strong>\1</strong>', text)

    # Convert code blocks
    text = re.sub(r'```(.*?)```', r'<pre><code>\1</code></pre>', text, flags=re.DOTALL)
    text = re.sub(r'`([^`]+?)`', r'<code>\1</code>', text)

    # Convert bullet points
    text = re.sub(r'\n\s*\* (.*?)\n', r'\n<ul><li>\1</li></ul>\n', text)
    text = re.sub(r'(<ul><li>.*?</li></ul>)', r'\1', text, flags=re.DOTALL)

    # Convert paragraphs
    text = re.sub(r'\n(?!<h|<ul|<pre|<code)(.*?)\n', r'\n<p>\1</p>\n', text)

    return text

def ViewSummary():
    with open('Data.json', 'r') as file:
        data = json.load(file)
    
    # Get the result text
    result_text = data.get('summarized_text', '')
    
    # Convert text to HTML
    content = convert_text_to_html(result_text)
    # Add basic HTML structure
    php_code = f"""
<?php 
session_start();
if(isset($_SESSION['username']) && isset($_SESSION['fullname']) && $_SESSION['role'] == 'student')
{{
?>
<!DOCTYPE html>
<html lang="en">
  <head>
    <meta charset="UTF-8" />
    <meta name="viewport" content="width=device-width, initial-scale=1.0" />
    <title>Summary View</title>
    <script src="https://cdnjs.cloudflare.com/ajax/libs/jspdf/2.4.0/jspdf.umd.min.js"></script>
    <script src="https://cdnjs.cloudflare.com/ajax/libs/html2canvas/1.4.1/html2canvas.min.js"></script>
    <link rel="stylesheet" href="../CSS/styles.css" />
    <style>
      body {{
        font-family: Arial, sans-serif;
        margin: 0;
        padding: 0;
      }}
      h1, h2 {{
          font-size: 1.5em;
                color: #0066cc;
                margin-bottom: 20px;
            }}
            p {{
                font-size: 18px;
                font-family: Arial, sans-serif;
                font-weight: 600;
                margin-bottom: 15px;
            }}
            pre, code {{
                font-family: 'Courier New', Courier, monospace;
                background-color: #f0f0f0;
                padding: 8px;
                border-radius: 5px;
                font-size: 0.9em;
                margin-bottom: 15px;
                overflow-x: auto; /* Enable horizontal scroll */
            }}
            ul {{
                list-style-type: disc;
                margin-bottom: 15px;
                padding-left: 20px;
            }}
            li {{
                margin-bottom: 5px;
            }}
            strong {{
                font-size: 1.4em;
                padding-top:20px;
                color: #009900;
                font-weight: bold;
            }}
      #content {{
        margin: 15px 25px 5px 25px;}}
      .container {{
        background: #fff;
        padding: 20px;
        border-radius: 12px;
        box-shadow: 0 0 10px rgba(0, 0, 0, 0.1);
        max-width: 1150px;
        width: 100%;
        margin: 20px auto;
        min-height: 525px;
      }}
      #content {{
        margin-bottom: 20px;
      }}
      #content h2 {{
        color: #28a745;
        font-size: 24px;
        margin-bottom: 10px;
      }}
      #content p {{
        color: #555;
        margin-bottom: 8px;
        font-size: 16px;
        line-height: 1.5;
      }}
      .btns {{
        display: flex;
        flex-direction: column;
        align-items: center;
        margin-bottom: 20px;
      }}
      .Download-button,
      .courses-button {{
        padding: 10px 20px;
        border: none;
        border-radius: 5px;
        cursor: pointer;
        font-size: 16px;
        display: block;
        margin-bottom: 5px;
        background-color: #320943e3;
        color: white;
        width: 200px;
      }}
      .courses-button:hover,
      .Download-button:hover {{
        background-color: #d6d6de;
        color: #000;
      }}
      .search-container {{
        position: absolute;
        display: flex;
        flex-direction: column;
        align-items: center;
        right: 280px;
        top: 31px;
      }}
      #search-icon {{
        cursor: pointer;
        width: 40px;
        border-radius: 50%;
      }}
      #search-input {{
        width: 500px;
        padding: 10px;
        border: 1px solid #d1d1d1;
        border-radius: 10px;
        display: none;
        position: absolute;
        top: 76px;
        background: #d1d1d1;
        box-shadow: #8b8b8b, -10px -10px 20px #ffffff;
        font-size: large;
      }}
      .highlight {{
        background-color: yellow;
      }}
      .loading-overlay {{
        position: fixed;
        top: 0;
        left: 0;
        width: 100%;
        height: 100%;
        background: rgba(255, 255, 255, 0.8);
        display: flex;
        justify-content: center;
        align-items: center;
        z-index: 9999;
        display: none;
      }}
      .spinner {{
        border: 4px solid rgba(0, 0, 0, 0.1);
        width: 36px;
        height: 36px;
        border-radius: 50%;
        border-top-color: #28a745;
        animation: spin 1s ease-in-out infinite;
      }}
      @keyframes spin {{
        to {{
          transform: rotate(360deg);
        }}
      }}
    </style>
  </head>
  <body>
    <nav>
      <div class="logo2">
        <img src="../images/logo.png" class="logo">
        <div class="logo-text">
          <p>Student Assistant</p>
        </div>
      </div>
      <ul>
        <li><a href="dashboard.php">Dashboard</a></li>
        <li><a href="course.php">Courses</a></li>
        <li><a href="quiz.php">Quizzes</a></li>
        <li><a href="summary.php" class="active">Summary</a></li>
        <li><a href="GPA.php">GPA Calculator</a></li>
      </ul>
      <div class="notification" onclick="toggleNotifi()">
        <img src="../images/notification.jpg" alt=""> <span>2</span>
      </div>
      <div class="notifi-box" id="notificationBox">
        <div class="notifi-item">
          <img src="../images/profile.png">
          <div class="text">
            <h4>Student Assistant</h4>
            <p>Feb 12, 2022</p>
            <p>Welcome to Student Assistant, we wish you all the best on your journey!</p>
          </div>
        </div>
        <div class="notifi-item" onclick="openModal()">
          <img src="../images/profile.png">
          <div class="text">
            <h4>Ali Ahmed</h4>
            <p>Feb 12, 2022</p>
          </div>
        </div>
      </div>
      <img src="../images/profile.png" onclick="ToggleMenu()" class="user-pic">
      <div class="sub-menu-wrap" id="subMenu">
        <div class="sub-menu">
          <div class="user-info">
            <?php if(isset($_SESSION['profile_image'])): ?>
              <img src="data:image/jpeg;base64,<?php echo $_SESSION['profile_image']; ?>">
            <?php endif; ?>
            <h3><?php echo $_SESSION['fullname']; ?></h3>
          </div>
          <hr>
          <a href="profile.php" class="sub-menu-link">
            <img src="../images/edit.png">
            <p>Profile</p>
            <span>></span>
          </a>
          <hr>
          <a class="sub-menu-link" id="dark-mode-toggle" style="cursor: pointer;">
            <img src="../images/swich.jpg">
            <p>Swich</p>
          </a>
          <hr>
          <a href="../php/logout.php" class="sub-menu-link">
            <img src="../images/logout.png">
            <p>Logout</p>
            <span>></span>
          </a>
        </div>
      </div>
      <div id="myModal" class="modal">
        <div class="modal-content">
          <span class="close">&times;</span>
          <h4>Ail Ahmed</h4>
          <p>Date: 12-12-2023</p>
          <p>Hi, Can i fix you?</p>
          <textarea id="responseMessage" placeholder="Write your response here..."></textarea>
          <button id="sendResponseBtn">Send Response</button>
          <p id="successMessage" class="success-message">Response sent successfully!</p>
        </div>
      </div>
    </nav>
    <div class="container" id="capture">
      <div id="content">
        {content}
      </div>
    </div>
    <div class="btns">
      <button class="Download-button" onclick="downloadPDF()">Download</button>
      <button class="courses-button" onclick="location.href='summary.php';">Back to Summary</button>
    </div>
    <div id="loading" class="loading-overlay">
      <div class="spinner"></div>
    </div>
    <footer>
      <div class="row">
        <div class="col">
          <img src="../images/logo.png" class="logo-footer">
          <h1>Student Assistant</h1>
          <h4>From study materials to assessments, we've got you covered.</h4>
          <div class="social">
            <i class="fa fa-facebook-f"></i>
            <i class="fa fa-twitter"></i>
            <i class="fa fa-instagram"></i>
            <i class="fa fa-linkedin"></i>
          </div>
        </div>
        <div class="col">
          <h3>Product</h3><br>
          <div class="line">
            <a href="course.php">Course management</a>
            <br><br>
            <a href="quiz.php">Quiz management</a>
            <br><br>
            <a href="summary.php">Summary</a>
            <br><br>
            <a href="GPA.php">GPA Calculator</a>
          </div>
        </div>
        <div class="col">
          <h3>Company</h3><br>
          <div class="line">
            <a href="aboutUs.php">About Us</a>
          </div>
        </div>
        <div class="col">
          <h3>Resources</h3><br>
          <div class="line">
            <a href="contactUs.php">Contact Us</a>
          </div>
        </div>
        <div class="copy">
          <p>Copyright ©2023-2024</p>
        </div>
      </div>
    </footer>
    <script src="../JS/main.js"></script>
    <script>
      function downloadPDF() {{
        const {{ jsPDF }} = window.jspdf;
        const contentElement = document.getElementById("capture");
        const loadingOverlay = document.getElementById("loading");

        loadingOverlay.style.display = "flex";

        html2canvas(contentElement, {{
          scale: 2,
        }}).then((canvas) => {{
          const imgData = canvas.toDataURL("image/png");
          const pdf = new jsPDF("p", "mm", "a4");
          const pdfWidth = pdf.internal.pageSize.getWidth();
          const pdfHeight = pdf.internal.pageSize.getHeight();
          const imgProps = pdf.getImageProperties(imgData);
          const imgWidth = imgProps.width;
          const imgHeight = imgProps.height;

          const ratio = imgWidth / imgHeight;
          const newHeight = pdfWidth / ratio;

          if (newHeight > pdfHeight) {{
            const newWidth = pdfHeight * ratio;
            pdf.addImage(
              imgData,
              "PNG",
              (pdfWidth - newWidth) / 2,
              0,
              newWidth,
              pdfHeight
            );
          }} else {{
            pdf.addImage(imgData, "PNG", 0, 0, pdfWidth, newHeight);
          }}

          pdf.save("download.pdf");
          loadingOverlay.style.display = "none";
        }}).catch(() => {{
          loadingOverlay.style.display = "none";
        }});
      }}

      document.addEventListener("DOMContentLoaded", function () {{
        const loadingOverlay = document.getElementById("loading");
        loadingOverlay.style.display = "flex";

        // Simulated JSON data for demonstration
        const jsonData = {{
          summarized_text: `
            **Heading 1**
            This is a paragraph of content.

            *Bullet point 1*
            *Bullet point 2*
          `
        }};

        const contentElement = document.getElementById("content");
        const lines = jsonData.summarized_text.split("\\n");

        lines.forEach((line) => {{
          if (line.startsWith("**")) {{
            const header = document.createElement("h2");
            header.textContent = line.replace(/\\*\\*/g, "");
            contentElement.appendChild(header);
          }} else if (line.startsWith("*")) {{
            const paragraph = document.createElement("p");
            paragraph.textContent = line.replace(/\\*/g, "");
            contentElement.appendChild(paragraph);
          }} else if (line.trim() === "") {{
            const br = document.createElement("br");
            contentElement.appendChild(br);
          }}
        }});

        loadingOverlay.style.display = "none";
      }});

      document.getElementById("search-input").addEventListener("input", function () {{
        const searchTerm = this.value.trim();
        const contentElement = document.getElementById("content");

        const highlightedElements = contentElement.querySelectorAll(".highlight");
        highlightedElements.forEach(element => {{
          element.classList.remove("highlight");
          element.innerHTML = element.textContent;
        }});

        if (searchTerm !== "") {{
          const regex = new RegExp(`({{searchTerm}})`, "gi");

          const textNodes = getTextNodes(contentElement);
          let firstMatchFound = false;

          textNodes.forEach(node => {{
            const matches = node.textContent.match(regex);
            if (matches) {{
              const newNode = document.createElement("span");
              newNode.innerHTML = node.textContent.replace(regex, '<span class="highlight">$1</span>');
              node.parentNode.replaceChild(newNode, node);

              if (!firstMatchFound) {{
                newNode.scrollIntoView({{ behavior: "smooth", block: "center" }});
                firstMatchFound = true;
              }}
            }}
          }});
        }}
      }});

      function getTextNodes(element) {{
        const textNodes = [];
        function recursiveSearch(node) {{
          if (node.nodeType === Node.TEXT_NODE) {{
            textNodes.push(node);
          }} else if (node.nodeType === Node.ELEMENT_NODE) {{
            node.childNodes.forEach(child => recursiveSearch(child));
          }}
        }}
        recursiveSearch(element);
        return textNodes;
      }}

      window.toggleSearchBar = function () {{
        const searchInput = document.getElementById("search-input");
        if (searchInput.style.display === "none" || searchInput.style.display === "") {{
          searchInput.style.display = "block";
        }} else {{
          searchInput.style.display = "none";
        }}
      }};
    </script>
  </body>
</html>
<?php 
 
}} else {{
  header("Location: index.php");
  exit();
}} 
?>
"""


    # Save to an HTML file
    with open('PAGES/summaryView.php', 'w') as file:
        file.write(php_code)

def save_to_database(user_id,summary_name,AIResponce): 
    conn = mysql.connector.connect(
        host="localhost", user="root", password="", database="student_assistance_db"
    )
    cursor = conn.cursor()
    sql = "INSERT INTO summary (student_id, summary_name, content) VALUES (%s, %s, %s)"

    val = (user_id, summary_name, AIResponce)
    cursor.execute(sql, val)
    conn.commit()
    print("Data saved to database successfully.")
    cursor.close()
    conn.close()


audio_wav = "temp_audio.wav"
audio_m4a = "temp_audio.m4a"

generation_config = {
    "temperature": 0.9,
    "top_p": 1,
    "top_k": 1,
    "max_output_tokens": 30000,
}


def process_conversion(video):
    print("Video converting...")
    command = [
        "ffmpeg",
        "-y",
        "-i",
        video,
        "-vn",
        "-ar",
        "16000",
        "-ac",
        "1",
        "-sample_fmt",
        "s16",
        audio_wav,
    ]

    try:
        subprocess.run(command, check=True)
        print("Conversion successful")
    except subprocess.CalledProcessError as e:
        print(f"Error during conversion: {e}")

    return process_compression_wav_to_m4a("temp_audio.wav")


def process_compression_wav_to_m4a(input_wav):
    print("Compressing WAV to M4A...")
    command = [
        "ffmpeg",
        "-i",
        input_wav,
        "-c:a",
        "aac",
        "-b:a",
        "192k",
        "-threads",
        "6",
        audio_m4a,
    ]

    try:
        if os.path.exists('temp_audio.m4a'):
        # File exists, remove it (force overwrite)
            os.remove('temp_audio.m4a')
            print("Cache cleared.")
        subprocess.run(command, check=True)
        print("Compression successful")
    except subprocess.CalledProcessError as e:
        print(f"Error during compression: {e}")

    if audio_m4a:
        return audio_m4a
    else:
        return "Compression failed"


def process_transcribe(audio):
    print("Audio Transcribing...")
    aai.settings.api_key = "6b25378de7044a129d97873d58f0183c"
    transcriber = aai.Transcriber()
    transcript = transcriber.transcribe(audio)
    tText = transcript.text
    transcript = {"Transcript": transcript.text}
    return tText


def proccess_gemini(prompt):
    genai.configure(api_key="AIzaSyDdpMjxVnjmaAKT1Bv_9rcow20bcUqHQLc")  

    generation_config = {
        "temperature": 1,
        "top_p": 0.95,
        "top_k": 64,
        "max_output_tokens": 8192,
        "response_mime_type": "text/plain",
    }

    model = genai.GenerativeModel(
        model_name="gemini-1.5-flash",
        generation_config=generation_config,
        system_instruction="Read the following content and provide a detailed summary with Minimum 1K word. The summary should be thorough yet concise, capturing all the key points and main ideas. Use simple and easy-to-understand language to ensure clarity. Aim for an extended summary that covers all the essential aspects without being overly verbose. The goal is to produce a summary that is both informative and accessible. just jump into the summarization proccess don't say anything at first.",
    )
    print("Starting chat session...")
    chat_session = model.start_chat(history=[])
    print("Sending message...")
    response = chat_session.send_message(prompt)
    return response.text

def process_summarize(originalText):
    print("Process Summarization...")

    return proccess_gemini(originalText)


def process_fetch(link):
    ydl_opts = {
        "outtmpl": "processedVideo.mp4",
        "format": "worstvideo[ext=mp4]+worstaudio[ext=m4a]/mp4",
        "verbose": True,
    }
    with youtube_dl.YoutubeDL(ydl_opts) as ydl:
        ydl.download([link])
    return "processedVideo.mp4"


def extract_text_from_pdf(file_path):
    text = ""
    with open(file_path, "rb") as file:
        reader = PdfReader(file)
        for page in reader.pages:
            text += page.extract_text()
    return text


def extract_text_from_ppt(file_path):
    text = ""
    presentation = Presentation(file_path)
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text


def extract_text_from_docx(file_path):
    text = ""
    doc = Document(file_path)
    for para in doc.paragraphs:
        text += para.text + "\n"
    return text


def process_request(file, link):
    global extracted_audio
    if file:
        filename = secure_filename(file.filename)
        file.save(filename)

        if filename.endswith(".pdf"):
            extracted_text = extract_text_from_pdf(filename)
        elif filename.endswith(".ppt") or filename.endswith(".pptx"):
            extracted_text = extract_text_from_ppt(filename)
        elif filename.endswith(".docx"):
            extracted_text = extract_text_from_docx(filename)
        else:
            global video
            video = filename

            extracted_audio = process_conversion(video)

            if extracted_audio:
                file_size = os.path.getsize(extracted_audio)
                print(f"Size of the audio file: {float(file_size/1024/1024)} MB")

                result_text = process_transcribe(extracted_audio)
                extracted_text = result_text if result_text else None
            else:
                return "Failed to fetch the audio!", 400

    elif link:
        video = process_fetch(link)
        extracted_audio = process_conversion(video)

        if extracted_audio:
            file_size = os.path.getsize(extracted_audio)
            print(f"Size of the audio file: {int(file_size/1024/1024)} MB")

            result_text = process_transcribe(extracted_audio)
            extracted_text = result_text if result_text else None
        else:
            return "Failed to fetch the audio!", 400
    else:
        return "No file or link found in request", 400

    if extracted_text is not None:
        summarized_text = process_summarize(extracted_text)
        global AIResponce
        AIResponce = summarized_text
        Data = {"summarized_text": summarized_text}
        with open("Data.json", "w", encoding="utf-8") as json_file:
            json.dump(Data, json_file, ensure_ascii=False, indent=2)
        return jsonify({"summarized_text": summarized_text})
    else:
        return "No response from the server.", 500


app = Flask(__name__)
CORS(app)


@app.route("/processing", methods=["POST"])
def Proccess():
    
    if "file" in request.files:
        uploaded_file = request.files["file"]
        content = process_request(file=uploaded_file, link=None)
        userId = request.form["user_id"]
        summary_name = request.form['summary_name']
        save_to_database(userId, summary_name,AIResponce)
        ViewSummary()
        return content
    elif "link" in request.form:
        video_link = request.form["link"]
        content = process_request(link=video_link, file=None)
        userId = request.form["user_id"]
        summary_name = request.form['summary_name']
        save_to_database(userId, summary_name,AIResponce)
        os.remove(extracted_audio)
        os.remove(audio_wav)
        os.remove(video)
        ViewSummary()
        return content
    else:
        return "No file or link found in request", 400


if __name__ == "__main__":
    app.run(debug=True, port=5055)
