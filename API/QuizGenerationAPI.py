from flask import Flask, jsonify, request
import json
from flask_cors import CORS
import google.generativeai as genai
import mysql.connector
import os
import time
from werkzeug.utils import secure_filename
from PyPDF2 import PdfReader
from pptx import Presentation
from docx import Document

app = Flask(__name__)
CORS(app)

global mcq_request
global tf_request
global cooldown_period

mcq_request = 0
tf_request = 0
cooldown_period = 60

def save_to_database(user_id,quiz_name,quizdata): 
    conn = mysql.connector.connect(
        host="localhost", user="root", password="", database="student_assistance_db"
    )
    cursor = conn.cursor()
    sql = "INSERT INTO quiz (student_id, quizname, quizdata, date_taken) VALUES (%s, %s, %s, now())"
    val = (user_id, quiz_name, quizdata)
    cursor.execute(sql, val)
    conn.commit()
    print("Data saved to database successfully.")
    cursor.close()
    conn.close()
    
    
def FetchOneMCQ(prompt):
    global mcq_request
    mcq_request = mcq_request + 1
    global cooldown_period

    if mcq_request == 15:
        print("The API is currently experiencing high traffic! Waiting for cooldown...")
        time.sleep(cooldown_period)
        mcq_request = 0

    genai.configure(api_key="AIzaSyDdpMjxVnjmaAKT1Bv_9rcow20bcUqHQLc")

    generation_config = {
        "temperature": 1,
        "top_p": 0.95,
        "top_k": 64,
        "max_output_tokens": 8192,
        "response_mime_type": "text/plain",
    }

    model = genai.GenerativeModel(
        model_name="gemini-1.5-flash",
        generation_config=generation_config,
        system_instruction="Using the content provided, generate one multiple-choice question (MCQ) along with four possible answers (A, B, C, D). Ensure that only one of the answers is correct. The question should be clear, concise, and relevant to the content. The last line will be the correct answer character (A, B, C, or D).\n\nThe format:\n*Question*\n*A. Choice 1*\n*B. Choice 2*\n*C. Choice 3*\n*D. Choice 4*\n*Correct Answer Character*",
    )
    chat_session = model.start_chat(history=[])
    response = chat_session.send_message(prompt)
    return response.text


def FetchOneTF(prompt):
    global tf_request
    tf_request = tf_request + 1
    global cooldown_period

    if tf_request == 15:
        print("The API is currently experiencing high traffic! Waiting for cooldown...")
        time.sleep(cooldown_period)
        tf_request = 0

    genai.configure(api_key="AIzaSyDdpMjxVnjmaAKT1Bv_9rcow20bcUqHQLc")

    generation_config = {
        "temperature": 1,
        "top_p": 0.95,
        "top_k": 64,
        "max_output_tokens": 8192,
        "response_mime_type": "text/plain",
    }

    model = genai.GenerativeModel(
        model_name="gemini-1.5-flash",
        generation_config=generation_config,
        system_instruction="Using the content provided, generate one True or False question. Ensure that the statement is clear, precise, and directly related to the content. Indicate whether the statement is true or false.\nThe last line will be the answer (True or False).\nThe format:\n*Question*\n*Correct Answer*",
    )
    chat_session = model.start_chat(history=[])
    response = chat_session.send_message(prompt)
    return response.text

def extract_text_from_pdf(file_path):
    text = ""
    with open(file_path, "rb") as file:
        reader = PdfReader(file)
        for page in reader.pages:
            text += page.extract_text()
    return text


def extract_text_from_ppt(file_path):
    text = ""
    presentation = Presentation(file_path)
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text


def extract_text_from_docx(file_path):
    text = ""
    doc = Document(file_path)
    for para in doc.paragraphs:
        text += para.text + "\n"
    return text

def process_request(file):
    filename = secure_filename(file.filename)
    file.save(filename)
    if filename.endswith(".pdf"):
        extracted_text = extract_text_from_pdf(filename)
    elif filename.endswith(".ppt") or filename.endswith(".pptx"):
        extracted_text = extract_text_from_ppt(filename)
    elif filename.endswith(".docx"):
        extracted_text = extract_text_from_docx(filename)
    os.remove(filename)
    return extracted_text

@app.route("/GeneratingQuiz", methods=["GET", "POST"])
def Main():
    if "Data"  in request.form:
        Data = request.form["Data"]
        contentData = json.loads(Data)["summarized_text"]
    elif "file" in request.files:
        file = request.files["file"]
        content = process_request(file)
        contentData = content
    else:
        return jsonify({"error": "No data provided!"}), 400

    userID = request.form["User_ID"]
    quizName = request.form["Quiz_name"]

    mcq_questions = []
    mcq_answers = []
    for _ in range(5):
        while True:
            question = FetchOneMCQ(contentData)
            lines = question.splitlines()
            question_text = "\n".join(lines[:-1]).strip()
            answer = lines[-1].strip()

            if (
                any(
                    answer.startswith(x)
                    for x in [
                        "A.",
                        "B.",
                        "C.",
                        "D.",
                        "a.",
                        "b.",
                        "c.",
                        "d.",
                        "A)",
                        "B)",
                        "C)",
                        "D)",
                        "a)",
                        "b)",
                        "c)",
                        "d)",
                        "**A.**",
                        "**B.**",
                        "**C.**",
                        "**D.**",
                        "*A.*",
                        "*B.*",
                        "*C.*",
                        "*D.*",
                        "**A**",
                        "**B**",
                        "**C**",
                        "**D**",
                        "*A*",
                        "*B*",
                        "*C*",
                        "*D*",
                        "A",
                        "B",
                        "C",
                        "D",
                    ]
                ) or "?\nA." in question_text or "?*\n*A." in question_text or "?\n\nA." in question_text
                 
            ):
                if "?\nA." in question_text:
                    question_text = question_text.replace("?\nA.", "?\n\nA.")
                    mcq_questions.append(question_text)
                    answer = answer.replace("*", "").replace("**", "")
                    mcq_answers.append(answer[0])
                    print("MCQ question added...")
                    break
                if "?*\n*A." in question_text:
                    question_text = question_text.replace("*", "")
                    question_text = question_text.replace("?\nA.", "?\n\nA.")
                    mcq_questions.append(question_text)
                    answer = answer.replace("*", "").replace("**", "")
                    mcq_answers.append(answer[0])
                    print("MCQ question added...")
                    break
                if "?\n\nA." in question_text:
                    mcq_questions.append(question_text)
                    answer = answer.replace("*", "").replace("**", "")
                    mcq_answers.append(answer[0])
                    print("MCQ question added...")
                    break
            else:
                print("Fixing Format...")

    tf_questions = []
    tf_answers = []
    for _ in range(5):
        while True:
            question = FetchOneTF(contentData)
            lines = question.splitlines()
            question_text = "\n".join(lines[:-1]).strip()
            answer = lines[-1].strip()

            if not any(line.strip() == "" for line in lines) and answer.lower() in (
                "true",
                "false",
                "**true**",
                "**false**",
                "*true*",
                "*false*",
            ):
                question_text = question_text.replace("*", "").replace("**", "")
                tf_questions.append(question_text)
                answer = answer.replace("*", "").replace("**", "")
                tf_answers.append(answer)
                print("T\\F question added...")

                break
            else:
                print(f"Fixing Format...")

    
    quiz_data = json.dumps({"Questions": mcq_questions + tf_questions, "Answers": mcq_answers + tf_answers})
    save_to_database(userID,quizName,quiz_data)
    with open("Quiz_data.json", "w", encoding="utf-8") as f:
        json.dump(
            quiz_data,
            f, ensure_ascii=False, indent=2
        )
        return (
        jsonify(
            {
                "Questions": mcq_questions + tf_questions,
                "Answers": mcq_answers + tf_answers,
            }
        ),
        200,
    )


if __name__ == "__main__":
    app.run(debug=False, port=5005)
